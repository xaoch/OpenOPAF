import sys
import csv
import os
from pptx import Presentation
import argparse
import csv
from threading import Thread, Event
import subprocess

class SlidesExtractor:

    stopSignal = Event()

    def __init__(self, directory):
        self.prs = Presentation(directory+"/presentation.pptx")
        self.path=directory+"/Slides"
        self.pathPres=directory+"/presentation.pptx"
        self.outputPDF=directory+"/presentation.pdf"
        self.outputSlides=self.path+"/slide%02d.png"
        self.createFolders()

        # you should use full paths, to make sure PowerPoint can handle the paths
        png_folder = ...
        pptx_file = ...

        self.csv_file = open(self.path + '/result.csv', mode='w')
        self.resultFile = csv.writer(self.csv_file, delimiter=',')
        self.resultFile.writerow(["slide", "font_size","text_length"])
        self.fs = 44100  # Sample rate
        self.interval = 5  # Duration of recording seconds

        self.minimum_pitch = 80
        self.maximum_pitch = 400
        self.time_step = 0.01
        print("Slides Thread Ready")

    def createFolders(self):
        print(self.path)
        try:
            os.mkdir(self.path)
        except:
            print("Slides Directories already created")

    def extract(self):
        print("Slides Thread: starting")
        fontErrors = [None] * len(self.prs.slides)
        textErrors = [None] * len(self.prs.slides)
        slides = 0
        for slide in self.prs.slides:
            slides = slides + 1
            shapes = 0
            textSlide = ""
            textlength = "No Text"
            fontsize = "No Font"
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                shapes = shapes + 1
                paragraphs = 0
                fontsize = "Ok"
                for paragraph in shape.text_frame.paragraphs:
                    paragraphs = paragraphs + 1
                    runs = 0
                    for run in paragraph.runs:
                        textSlide=textSlide+ " "+ run.text
                        runs = runs + 1
                        size = run.font.size
                        if size is not None:
                            size=size.pt
                        else:
                            size=18
                        if size is not None:
                            if size < 18:
                                 fontsize = "Small"
                            if size < 14:
                                fontsize = "Tiny"
            textlength="Ok"
            if len(textSlide.split()) > 60:
                textlength="Long"
            elif len(textSlide.split()) > 33:
                textlength="Verbose"
            fontErrors[slides - 1] = fontsize
            textErrors[slides-1] = textlength
        print(fontErrors)
        print(textErrors)
        index = 0

        for i in range(slides):
            index = i + 1
            self.resultFile.writerow([index, fontErrors[i], textErrors[i]])
        self.csv_file.close()
        command=["libreoffice","--headless","--convert-to","pdf",self.pathPres]
        subprocess.run(command)
        command=["gs","-sDEVICE=pngalpha","-o",self.outputSlides,self.outputPDF]
        subprocess.run(command)
        print("Slides Thread: finishing")

    def run(self):
        self.t = Thread(target=self.extract, args=())
        self.t.start()

    def stop(self):
        self.stopSignal.set()
        self.t.join()





